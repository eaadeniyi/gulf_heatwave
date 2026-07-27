"""
Build a presentation-ready APPENDIX deck (.pptx) summarizing the heatwave
project's definitions, data dictionary, methods, naming rationale, and results.
Content mirrors REFERENCE_glossary_methods_results.md. python-pptx only.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Heatwave_Reference_Appendix.pptx")

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xC4, 0x4E, 0x52)
HEADER_BG = RGBColor(0x1F, 0x3A, 0x5F)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ALT_BG = RGBColor(0xEF, 0xF1, 0xF5)
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def title_bar(slide, title, subtitle=None):
    tb = _tb(slide, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    if subtitle:
        p2 = tb.text_frame.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(12); r2.font.color.rgb = GREY


def footnote(slide, text):
    tb = _tb(slide, Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(8); r.font.italic = True; r.font.color.rgb = GREY


def add_title_slide():
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, Inches(2.4), SW, Inches(1.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Heatwave Classification — Methods, Data Dictionary & Key Results"
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = HEADER_FG
    st = _tb(s, Inches(0.65), Inches(4.3), Inches(12), Inches(2))
    for line, sz in [("Appendix / reference for the presentation", 16),
                     ("Texas county-level heatwave classification  |  study period 2015–2025  |  climate baseline from 1979", 12),
                     ("Descriptive exposure classification only — not injury, worker heat dose, or official NWS advisories", 11)]:
        p = st.text_frame.add_paragraph() if st.text_frame.paragraphs[0].runs else st.text_frame.paragraphs[0]
        rr = p.add_run(); rr.text = line; rr.font.size = Pt(sz)
        rr.font.color.rgb = GREY if sz < 16 else NAVY


def add_section(title):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, Inches(3.0), SW, Inches(1.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = HEADER_FG


def add_bullets(title, bullets, subtitle=None, note=None):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title, subtitle)
    body = _tb(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4))
    tf = body.text_frame
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        r = p.add_run(); r.text = ("• " if level == 0 else "– ") + text
        r.font.size = Pt(16 if level == 0 else 13)
        r.font.color.rgb = NAVY if level == 0 else GREY
        p.space_after = Pt(6)
    if note:
        footnote(s, note)


def add_table(title, header, rows, col_widths, subtitle=None, note=None, body_pt=11, max_rows=13):
    """Add one or more table slides (auto-paginates rows)."""
    chunks = [rows[i:i + max_rows] for i in range(0, len(rows), max_rows)] or [[]]
    for ci, chunk in enumerate(chunks):
        s = prs.slides.add_slide(BLANK)
        title_bar(s, title + (" (cont.)" if ci else ""), subtitle if ci == 0 else None)
        nrows = len(chunk) + 1
        ncols = len(header)
        total_w = Inches(12.3)
        left = Inches(0.5); top = Inches(1.45)
        height = Inches(min(5.3, 0.42 * nrows))
        gtable = s.shapes.add_table(nrows, ncols, left, top, total_w, height).table
        # column widths
        wsum = sum(col_widths)
        for j, w in enumerate(col_widths):
            gtable.columns[j].width = Emu(int(int(total_w) * w / wsum))
        # header
        for j, h in enumerate(header):
            c = gtable.cell(0, j)
            c.fill.solid(); c.fill.fore_color.rgb = HEADER_BG
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = h; r.font.size = Pt(body_pt + 1); r.font.bold = True; r.font.color.rgb = HEADER_FG
        # body
        for i, row in enumerate(chunk, start=1):
            for j, val in enumerate(row):
                c = gtable.cell(i, j)
                c.fill.solid(); c.fill.fore_color.rgb = ALT_BG if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = c.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                r = p.add_run(); r.text = str(val); r.font.size = Pt(body_pt)
                r.font.color.rgb = NAVY if j == 0 else RGBColor(0x22, 0x22, 0x22)
                if j == 0:
                    r.font.bold = True
        if note and ci == len(chunks) - 1:
            footnote(s, note)


# ======================================================================
# BUILD
# ======================================================================
add_title_slide()

add_bullets("What's in this appendix",
    [(0, "The two heatwave definitions and the shared design choices"),
     (0, "The heat-index metric and why it is called a 'proxy'"),
     (0, "How thresholds, heatwave days, and events are calculated"),
     (0, "Data dictionary — every column in every output table, with meaning"),
     (0, "Terminology & naming rationale (why we chose each name)"),
     (0, "Statistical computations used, and key results / conclusions"),
     (0, "Caveats & limitations to state alongside any result")],
    subtitle="A reference glossary for the methods-heavy parts of the analysis")

# ---- Definitions ----
add_section("1 · The heatwave definitions")
add_table("The two definitions",
          ["Definition", "Rule", "Plain meaning"],
          [["Definition 01", "county-relative 85th-pctl daily-mean heat index, >=2 consecutive days, walk-forward baseline",
            "days in the hottest ~15% for that county & time of year, sustained >=2 days"],
           ["Definition 02", "county-relative 95th-pctl daily-mean heat index, >=2 consecutive days, walk-forward baseline",
            "days in the hottest ~5% for that county & time of year, sustained >=2 days"]],
          col_widths=[2, 5.5, 4.5], body_pt=12,
          subtitle="Only the percentile changes between them; everything else is shared")
add_table("Shared design choices (both definitions)",
          ["Element", "Choice", "Why"],
          [["Metric", "daily-MEAN heat index", "the day's overall apparent-heat burden, not just the afternoon peak"],
           ["Baseline", "walk-forward (expanding): year Y uses 1979 … Y-1", "classify each year against all history observed up to that point"],
           ["Threshold", "county-relative percentile", "'unusually hot for THIS county and time of year'"],
           ["Persistence", ">= 2 consecutive days", "a heatwave is sustained heat, not one hot day"],
           ["Absolute floor", "none in primary (>=80F floor = sensitivity)", "faithful to the definition as written"],
           ["Artifacts", "3 confirmed RH-clip days set to missing", "bad data must not create heatwave days"],
           ["Windows", "centered 15-day (+/-7) AND calendar-month, reported side by side", "test sensitivity to how the percentile is pooled"]],
          col_widths=[2, 4.5, 5.5], body_pt=11)

# ---- Metric ----
add_section("2 · The heat metric (heat-index proxy)")
add_bullets("Heat index = 'feels-like' temperature",
    [(0, "Combines air temperature and relative humidity into one apparent-temperature value (NWS Rothfusz regression, with the official low/high-humidity adjustment terms)."),
     (0, "Two proxies are computed:"),
     (1, "daily-MEAN HI = heat_index(Tmean, mean RH)  →  used by the DEFINITIONS.  [column: derived_tmean_meanrh_hi_f]"),
     (1, "daily-MAX HI = heat_index(Tmax, RHmin)  →  used by the NWS proxy.  [column: derived_tmax_rhmin_hi_proxy_f]"),
     (0, "Why 'proxy', not 'heat index':"),
     (1, "inputs are DAILY (Tmax/Tmin/RHmax/RHmin), not hourly-concurrent — a derived approximation, not an observed hourly maximum"),
     (1, "temperature (GHCN stations) and humidity (gridMET grid) have DIFFERENT spatial support"),
     (0, "Data: temperature = NOAA GHCN-Daily; humidity = gridMET; geometry = 2020 Census TIGER/Line.")],
    subtitle="Tmean = (Tmax+Tmin)/2 ; mean RH = (RHmax+RHmin)/2")

# ---- Methods ----
add_section("3 · How thresholds & classification are computed")
add_bullets("From weather to heatwave events — 4 steps",
    [(0, "1. THRESHOLD (per county, per calendar slot, per year): take the 85th/95th percentile of baseline heat index in the window (±7 days, or the calendar month), pooled over 1979…Y-1. → threshold_value_f"),
     (0, "2. CANDIDATE day: daily-mean HI strictly > its own threshold.  → exceedance_f = mean-HI − threshold"),
     (0, "3. HEATWAVE day: a candidate day inside a run of >= 2 consecutive calendar days (a run breaks on a missing/non-consecutive/non-candidate day)."),
     (0, "4. HEATWAVE event: one uninterrupted run of heatwave days in one county; duration = end − start + 1."),
     (0, "Percentile via numpy.percentile (linear interpolation).")],
    subtitle="Walk-forward: the threshold is re-estimated every analysis year from all prior years")
add_bullets("IDW gap-filling of missing temperature (statewide)",
    [(0, "Many rural counties lack station data on some/all days. Missing daily temperature is filled by INVERSE-DISTANCE WEIGHTING from surrounding counties."),
     (0, "For a missing county-day:  value = Σ(wᵢ·vᵢ) / Σ(wᵢ)  over counties with data that day,"),
     (1, "weight wᵢ = 1 / (distance)²,  distance = between county CENTROIDS (equal-area projection, EPSG:5070)"),
     (0, "Every imputed county-day is FLAGGED (temp_imputed) so interpolated values never pass as observed."),
     (0, "Statewide: 12.8% of temperature county-days imputed; 22 counties had NO native station data; 93 were fully native.")],
    subtitle="Why IDW: gives every county a value while down-weighting distant counties (1/d²)")

# ---- Data dictionary ----
add_section("4 · Data dictionary (columns by table)")
add_table("Daily heatwave-day table — one row per heatwave county-date",
          ["Column", "Meaning / how computed"],
          [["county_fips / county_name", "county identity (FIPS key; name from Census)"],
           ["date / year / month", "calendar date of the heatwave day"],
           ["tmax_f / tmin_f / tmean_f", "daily max/min/mean air temperature (°F); Tmean=(Tmax+Tmin)/2"],
           ["rmin_pct", "daily minimum relative humidity (%), gridMET"],
           ["derived_tmean_meanrh_hi_f (hi_mean_f)", "daily-mean heat-index proxy = heat_index(Tmean, mean RH)"],
           ["threshold_value_f", "this county-date-year's percentile threshold (°F)"],
           ["exceedance_f", "mean-HI − threshold (°F above the bar)"],
           ["heatwave_day_flag", "1 = this county-date qualifies as a heatwave day"],
           ["event_id / event_duration_days", "the event this day belongs to, and that event's length"],
           ["temp_imputed", "True if the temperature was IDW gap-filled"],
           ["qc_status", "data-quality label (valid / suspicious_retain / missing_input / invalid_physical)"]],
          col_widths=[3.5, 8.5], body_pt=11)
add_table("Event table — one row per heatwave event",
          ["Column", "Meaning"],
          [["event_label", "readable id: <countyFIPS>_<onsetYear>_<seq> (e.g. 48201_2023_012)"],
           ["start_date / end_date", "first and last day of the event"],
           ["event_duration_days", "integer consecutive-day length"],
           ["peak_mean_hi_f", "highest daily-mean HI during the event (°F)"],
           ["peak_day_date", "the event's hottest day (by mean HI)"],
           ["peak_day_tmax_f / _tmean_f / _rmin_pct", "temperature & humidity on the peak day"],
           ["peak_day_threshold_f", "the threshold in force on the peak day"],
           ["tmax_max_f / tmean_mean_f", "hottest Tmax during event / mean of Tmean over event days"],
           ["peak_exceedance_f", "largest single-day exceedance (°F above threshold)"],
           ["cumulative_exceedance_f", "Σ positive exceedance over the event (°F·days) = anomaly 'dose'"],
           ["n_imputed_days / event_contains_imputed_day", "how many event days used IDW-filled temperature"],
           ["onset_year", "year the event started (used to count annual events)"]],
          col_widths=[3.7, 8.3], body_pt=10.5)
add_table("County-month and County-year summaries",
          ["Column", "Table", "Meaning"],
          [["heatwave_events_started", "both", "events whose ONSET is in this month / year"],
           ["heatwave_events_active", "county-month", "events overlapping this month (may have started earlier)"],
           ["heatwave_days", "both", "heatwave days in this month / year"],
           ["longest_event_duration_days", "both", "longest event touching this month / that year"],
           ["event_ids_started / _active", "county-month", "the actual event labels (started vs active)"],
           ["first_event_start_date / last_event_end_date", "county-year", "first-to-last event span"],
           ["heatwave_days_imputed", "county-year", "how many of the year's heatwave days used IDW temp"]],
          col_widths=[3.6, 2.2, 6.2], body_pt=11,
          note="Month-crossing events: counted once at onset month; 'active' in every month touched; DAYS allocated to their actual month (no double-counting).")
add_table("Threshold, quality-control & NWS-proxy columns",
          ["Column / value", "Meaning"],
          [["threshold_value_f / n_reference_values", "the percentile threshold and how many baseline obs fed it"],
           ["percentile / window_method", "85 or 95 ; 'centered 15-day' or 'calendar-month bucket'"],
           ["threshold_quality_flag", "'low_n_ref' if <20 baseline obs, else 'ok'"],
           ["qc_status", "valid / suspicious_retain / missing_input / invalid_physical"],
           ["qc_rh_pin_likely_artifact", "RH clipped to exactly 100% on a warm, rain-free day (bad data)"],
           ["rh_pin_class", "confirmed_artifact / likely_real_wet / indeterminate"],
           ["nws_office / advisory_hi_f / extreme_warning_hi_f", "assigned office + its advisory / extreme-warning HI thresholds"],
           ["nws_advisory_threshold_met", "1 if daily-max HI proxy >= advisory threshold"],
           ["advisory_threshold_days", "annual count of advisory-threshold days"],
           ["verification_status", "documented / sr_standard / approximate (honesty flag on thresholds)"]],
          col_widths=[4.2, 7.8], body_pt=10.5)

# ---- Terminology ----
add_section("5 · Terminology & naming rationale")
add_table("Why we chose each term / name",
          ["Term / choice", "Why we use it"],
          [["Heatwave day (county-date)", "unit of analysis is one county on one date; avoids the ambiguous 'event-day'"],
           ["Heatwave event (one run, one county)", "keeps each real, uninterrupted heat spell as its own record"],
           ["Event duration (integer days)", "reports real calendar length; we do NOT headline a pooled average like '4.3 days'"],
           ["'persistent apparent-heat anomaly'", "precise label for the year-round RELATIVE construct (unusual-for-the-date, not always absolutely hot)"],
           ["'proxy' (heat index, NWS)", "flags values derived from daily (not hourly) data; can't reproduce official products"],
           ["QA-only pooled totals", "cross-county sums are sanity checks, not the headline; substance is county-level"],
           ["self-describing columns", "threshold_value_f, exceedance_f, n_reference_values — readable without the code"],
           ["walk-forward", "names the expanding-baseline design (vs a fixed climatology)"]],
          col_widths=[3.8, 8.2], body_pt=11)

# ---- Stats ----
add_section("6 · Statistical computations & 7 · Results")
add_table("Statistical computations used",
          ["Computation", "What / why"],
          [["Percentile threshold (85th/95th)", "numpy.percentile of baseline HI; defines 'unusually hot for here'"],
           ["Walk-forward pooling", "baseline re-estimated yearly from 1979…Y-1; classify vs prior climate"],
           ["Persistence run-length", "consecutive-day counting with break rules; enforces the >=2-day rule"],
           ["Exceedance / cumulative exceedance", "HI − threshold ; Σ positive exceedance = event 'dose'"],
           ["IDW imputation", "1/distance² centroid-weighted fill of missing temperature (flagged)"],
           ["Jaccard overlap = |A∩B|/|A∪B|", "how much two definitions/methods agree on WHICH days"],
           ["Descriptive trend slope", "linear fit of annual days vs year — descriptive only, NOT formal trend inference"],
           ["Fixed-vs-walk-forward comparison", "does the expanding baseline suppress later-year counts / trend?"],
           ["Anchor-vs-composite sensitivity", "how sensitive are results to the multi-station composite temperature?"]],
          col_widths=[3.8, 8.2], body_pt=10.5)
add_table("Key result — Definition 01 vs 02 (statewide, 254 counties, 15-day window)",
          ["Metric", "Def 01 (85th)", "Def 02 (95th)"],
          [["pooled heatwave-days (QA)", "~170,900", "~52,800"],
           ["pooled events (QA)", "~48,300", "~17,400"],
           ["per-county heatwave days — MEDIAN", "677", "196"],
           ["per-county range", "154 – 1,230", "18 – 516"]],
          col_widths=[6, 3, 3], body_pt=13,
          note="The 95th-pctl definition yields ~1/3 the heatwave days of the 85th — it keeps only each county's most anomalous ~5% of days.")
add_bullets("Key results & conclusions",
    [(0, "ABSOLUTE-FLOOR sensitivity: adding a mean-HI>=80°F floor roughly HALVES the counts — a relative-only definition flags many sub-80°F COOL-SEASON anomaly days."),
     (1, "e.g. Harris Co. December 2021 (record-warm): ~19 heatwave days at 80–85°F → ~0 under the 80°F floor."),
     (0, "FIXED vs WALK-FORWARD baseline: day-level Jaccard 0.92; fixed flags +6.5% more days and steeper slopes → walk-forward absorbs part of the warming signal; rankings unchanged."),
     (0, "ANCHOR-STATION vs composite temperature: heatwave-day Jaccard only 0.45–0.73 → the TEMPERATURE SOURCE changes results more than the baseline choice. County-to-county map texture is unreliable; trust the regional gradient."),
     (0, "RELATIVE vs ABSOLUTE (NWS proxy): arid El Paso/Lubbock ≈1 advisory-threshold day in 11 yrs yet many relative heatwave days; humid Houston logs 172."),
     (0, "DATA-QUALITY: Mar 1 2023 gridMET clipped RH=100% in 118/254 counties, inflating HI +15–24°F → confirmed artifact, set to missing.")],
    subtitle="Section 7")

# ---- Caveats ----
add_section("8 · Caveats & limitations")
add_bullets("State these alongside any result",
    [(0, "Daily PROXY, not hourly apparent heat; temperature & humidity have different spatial support."),
     (0, "Composite-station + IDW NOISE: county temperature is a changing multi-station composite plus interpolation in low-coverage counties → single-county values are noisy; regional gradients are robust."),
     (0, "Year-round RELATIVE construct = 'persistent apparent-heat anomaly'; ~half the flagged days are sub-80°F cool-season anomalies unless the >=80°F floor is applied."),
     (0, "NWS proxy is APPROXIMATE (nearest-office crosswalk; most office thresholds flagged approximate) and is NOT an official advisory."),
     (0, "Trend slopes are DESCRIPTIVE (11 annual points) — not formal trend inference."),
     (0, "This is EXPOSURE classification only — not an injury or worker-heat-dose measure.")],
    note="Code (state-agnostic, config-driven) & data: github.com/eaadeniyi/gulf_heatwave  |  Def 01 = PERCENTILES [85], Def 02 = [95]")

prs.save(OUT)
print("[done] wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
