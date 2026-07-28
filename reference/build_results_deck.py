"""
Build the RESULTS deck (.pptx): per-definition results with embedded figures,
plain-language meaning, and a likely/scenario-questions section. python-pptx only.
Content mirrors RESULTS_PRESENTATION.md.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
OUT = os.path.join(ROOT, "texas_heatwave_pilot", "outputs", "TX")
FIG85 = os.path.join(OUT, "def_p85_2d", "figures")
FIG95 = os.path.join(OUT, "def_p95_2d", "figures")
OUTPPTX = os.path.join(HERE, "Heatwave_Results_Deck.pptx")

NAVY = RGBColor(0x1F, 0x3A, 0x5F); BLUE = RGBColor(0x4C, 0x72, 0xB0); RED = RGBColor(0xC4, 0x4E, 0x52)
HDR = RGBColor(0x1F, 0x3A, 0x5F); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ALT = RGBColor(0xEF, 0xF1, 0xF5); GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height


def tbx(s, l, t, w, h):
    b = s.shapes.add_textbox(l, t, w, h); b.text_frame.word_wrap = True; return b


def title_bar(s, title, sub=None, color=NAVY):
    b = tbx(s, Inches(0.5), Inches(0.25), Inches(12.3), Inches(1.0))
    p = b.text_frame.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = color
    if sub:
        p2 = b.text_frame.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(12); r2.font.color.rgb = GREY


def footnote(s, text):
    b = tbx(s, Inches(0.5), Inches(7.04), Inches(12.3), Inches(0.4))
    r = b.text_frame.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(8); r.font.italic = True; r.font.color.rgb = GREY


def title_slide():
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, Inches(2.3), SW, Inches(1.8)); bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Heatwave Definitions — Results, Figures & Meaning"
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
    b = tbx(s, Inches(0.65), Inches(4.35), Inches(12), Inches(2))
    for line, sz, col in [("Two definitions presented on their own — statewide Texas, 254 counties, 2015–2025", 15, NAVY),
                          ("Def 01 = county-relative 85th-pctl daily-mean heat index, ≥2 consecutive days, walk-forward baseline", 12, GREY),
                          ("Def 02 = same, at the 95th percentile", 12, GREY),
                          ("Descriptive exposure classification only — not injury, worker heat dose, or official NWS advisories", 11, RED)]:
        p = b.text_frame.add_paragraph() if b.text_frame.paragraphs[0].runs else b.text_frame.paragraphs[0]
        rr = p.add_run(); rr.text = line; rr.font.size = Pt(sz); rr.font.color.rgb = col


def section(title, color=NAVY):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, Inches(2.9), SW, Inches(1.5)); bar.fill.solid()
    bar.fill.fore_color.rgb = color; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE


def numbers_slide(title, rows, color, sub=None):
    s = prs.slides.add_slide(BLANK); title_bar(s, title, sub, color)
    n = len(rows) + 1
    t = s.shapes.add_table(n, 2, Inches(0.6), Inches(1.5), Inches(12.1), Inches(min(5.2, 0.44 * n))).table
    t.columns[0].width = Inches(7.2); t.columns[1].width = Inches(4.9)
    for j, h in enumerate(["Statistic", "Value"]):
        c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = HDR; c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = h; r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for i, (a, b) in enumerate(rows, 1):
        for j, val in enumerate([a, b]):
            c = t.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = ALT if i % 2 == 0 else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = str(val); r.font.size = Pt(12)
            r.font.color.rgb = NAVY if j == 0 else RGBColor(0x22, 0x22, 0x22); r.font.bold = (j == 1)
    return s


def image_slide(title, img, caption, color=NAVY):
    s = prs.slides.add_slide(BLANK); title_bar(s, title, color=color)
    im = Image.open(img); iw, ih = im.size; ar = iw / ih
    max_w, max_h = Inches(11.6), Inches(4.9)
    w = max_w; h = Emu(int(int(w) / ar))
    if int(h) > int(max_h):
        h = max_h; w = Emu(int(int(h) * ar))
    left = Emu(int((int(SW) - int(w)) / 2)); top = Inches(1.45)
    s.shapes.add_picture(img, left, top, width=w, height=h)
    b = tbx(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7))
    r = b.text_frame.paragraphs[0].add_run(); r.text = caption; r.font.size = Pt(12); r.font.color.rgb = GREY; r.font.italic = True


def bullets_slide(title, items, sub=None, color=NAVY, note=None):
    s = prs.slides.add_slide(BLANK); title_bar(s, title, sub, color)
    b = tbx(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(5.4)); tf = b.text_frame
    first = True
    for lvl, txt, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.level = lvl; r = p.add_run(); r.text = ("• " if lvl == 0 else "   – ") + txt
        r.font.size = Pt(14 if lvl == 0 else 12); r.font.bold = bold
        r.font.color.rgb = NAVY if lvl == 0 else GREY; p.space_after = Pt(5)
    if note:
        footnote(s, note)


def qa_slide(title, qas, color=NAVY):
    s = prs.slides.add_slide(BLANK); title_bar(s, title, color=color)
    b = tbx(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(5.5)); tf = b.text_frame
    first = True
    for q, a in qas:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = "Q: " + q; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY
        p.space_before = Pt(6)
        pa = tf.add_paragraph(); ra = pa.add_run(); ra.text = "A: " + a; ra.font.size = Pt(11); ra.font.color.rgb = GREY


# ================= BUILD =================
title_slide()
bullets_slide("What we are presenting",
    [(0, "Two heatwave definitions, each shown on its own (not compared head-to-head):", True),
     (1, "Definition 01 — county-relative 85th-percentile daily-mean heat index, >=2 consecutive days, walk-forward baseline", False),
     (1, "Definition 02 — same, at the 95th percentile", False),
     (0, "Unit of analysis:", True),
     (1, "heatwave DAY = one county on one date inside a >=2-day run", False),
     (1, "heatwave EVENT = one uninterrupted run of heatwave days in one county; duration = its length in days", False),
     (0, "Both are YEAR-ROUND, COUNTY-RELATIVE ('unusually hot for this county & date') = a persistent apparent-heat anomaly", True),
     (0, "For each definition: results from the CSVs -> figures -> what they mean; then likely questions; then caveats.", False)])

# ---- Definition 01 ----
section("Definition 01 — 85th percentile", BLUE)
numbers_slide("Definition 01 — results from the data", [
    ("Total heatwave county-days (statewide, 11 yr)", "170,894"),
    ("Total heatwave events", "48,323"),
    ("Heatwave days per county — median (p25–p75; range)", "677  (568–775; 154–1,230)"),
    ("Event duration — median / mean / max", "3 d / 3.5 d / 48 d"),
    ("Events that are the minimum 2 days", "44%"),
    ("Events lasting >= 5 days", "21%"),
    ("Share of heatwave days in Jun–Sep", "37%"),
    ("Top counties by heatwave days", "La Salle 1,230; Hudspeth 1,198; Lavaca 1,122"),
    ("Longest event", "Presidio, Jul 5–Aug 21 2023 = 48 days (peak mean-HI 99.5°F)"),
], BLUE, sub="Source: outputs/TX/def_p85_2d/ (centered 15-day window)")
image_slide("Def 01 — heatwave days per county (2015–2025)", os.path.join(FIG85, "map01_heatwave_days_per_county.png"),
            "A broad 'elevated-exposure' net: the typical county sees ~677 heatwave days over 11 years (~62/yr). Read the REGIONAL gradient, not single counties.", BLUE)
image_slide("Def 01 — when heatwave days fall", os.path.join(FIG85, "res_seasonal.png"),
            "Only ~37% of days are in summer (Jun–Sep, outlined). ~63% are cool-season 'unusual-for-the-date' anomalies — the signature of a year-round RELATIVE definition.", BLUE)
image_slide("Def 01 — how long events last", os.path.join(FIG85, "res_event_duration.png"),
            "Mostly short: 44% of events are the minimum 2 days (median 3). A minority are long persistent ridges (21% >=5 days; up to 48 days).", BLUE)

# ---- Definition 02 ----
section("Definition 02 — 95th percentile", RED)
numbers_slide("Definition 02 — results from the data", [
    ("Total heatwave county-days (statewide, 11 yr)", "52,786"),
    ("Total heatwave events", "17,428"),
    ("Heatwave days per county — median (p25–p75; range)", "196  (148–254; 18–516)"),
    ("Event duration — median / mean / max", "2 d / 3.0 d / 31 d"),
    ("Events that are the minimum 2 days", "53%"),
    ("Events lasting >= 5 days", "14%"),
    ("Share of heatwave days in Jun–Sep", "36%"),
    ("Top counties by heatwave days", "La Salle 516; Presidio 513; Hudspeth 508"),
    ("Longest event", "Kendall, Jul 2022 = 31 days (peak mean-HI 106.2°F)"),
], RED, sub="Source: outputs/TX/def_p95_2d/ (centered 15-day window)")
image_slide("Def 02 — heatwave days per county (2015–2025)", os.path.join(FIG95, "map01_heatwave_days_per_county.png"),
            "A severe-tail screen: the typical county sees ~196 heatwave days over 11 years (~18/yr) — each county's most anomalous ~5% of days.", RED)
image_slide("Def 02 — when heatwave days fall", os.path.join(FIG95, "res_seasonal.png"),
            "~36% of days in Jun–Sep — essentially the same seasonal spread as Def 01, confirming the cool-season share is a property of the RELATIVE method, not the percentile.", RED)
image_slide("Def 02 — how long events last", os.path.join(FIG95, "res_event_duration.png"),
            "Shorter & sparser than Def 01: 53% of events are exactly 2 days (median 2). Long events are rarer (14% >=5 days).", RED)

# ---- Threshold-window robustness ----
image_slide("Def 01 — heatwave days per county [calendar-month window]", os.path.join(FIG85, "map01_heatwave_days_per_county_month.png"),
            "Robustness: the calendar-month threshold gives essentially the same map as the 15-day window (per-county r=0.994).", BLUE)
image_slide("Def 01 — seasonal share [calendar-month window]", os.path.join(FIG85, "res_seasonal_month.png"),
            "Same seasonal pattern as the 15-day window (~37% of heatwave days in Jun-Sep).", BLUE)
image_slide("Def 02 — heatwave days per county [calendar-month window]", os.path.join(FIG95, "map01_heatwave_days_per_county_month.png"),
            "Robustness: near-identical to the 15-day window (per-county r=0.987).", RED)
image_slide("Def 02 — seasonal share [calendar-month window]", os.path.join(FIG95, "res_seasonal_month.png"),
            "Same seasonal pattern as the 15-day window (~36% of heatwave days in Jun-Sep).", RED)
numbers_slide("Both definitions were run on BOTH threshold windows", [
    ("Def 01 pooled days — 15-day / month", "170,894 / 171,115"),
    ("Def 01 pooled events — 15-day / month", "48,323 / 47,470"),
    ("Def 01 per-county median days — 15-day / month", "677 / 678   (r = 0.994)"),
    ("Def 02 pooled days — 15-day / month", "52,786 / 53,273"),
    ("Def 02 pooled events — 15-day / month", "17,428 / 17,517"),
    ("Def 02 per-county median days — 15-day / month", "196 / 194   (r = 0.987)"),
], NAVY, sub="Centered 15-day (primary, shown) vs calendar-month bucket — both computed for each definition")
s = prs.slides[-1]
footnote(s, "The two windows agree to <1% and r~=0.99 per county -> the window choice is a robustness check, not a driver. Month-window CSVs: *_month.csv in each def's tables/ folder.")

# ---- Companion ----
section("Companion outputs (definition-independent)")
bullets_slide("Data coverage & the NWS advisory-threshold proxy",
    [(0, "DATA COVERAGE / IDW imputation (same for both definitions):", True),
     (1, "~13% of temperature county-days are inverse-distance interpolated; 22 counties have NO native station; 93 are fully native", False),
     (1, "=> single-county map texture is noisy; the REGIONAL gradient is the trustworthy signal (see map03)", False),
     (0, "NWS advisory-threshold PROXY (uses the daily-MAX HI vs each county's local office threshold):", True),
     (1, "humid eastern/coastal counties accumulate many advisory-threshold days; arid far-west counties record almost none", False),
     (1, "=> a day can be a RELATIVE heatwave without reaching an ABSOLUTE advisory level (relative vs absolute heat)", False),
     (1, "PROXY only — not an official NWS advisory (no hourly data, no duration/overnight/coverage rules)", False)])

# ---- Q&A ----
section("Likely / scenario questions")
qa_slide("Likely questions — construct & metric", [
    ("Only ~37% of days are in summer — how is a warm January day a 'heatwave'?",
     "By design: a year-round, county-relative rule flags days unusual FOR THE DATE, not absolutely hot. ~63% are cool-season anomalies. Apply the optional mean-HI>=80°F floor (roughly halves counts) or a warm-season restriction for an absolute-heat framing."),
    ("Why a daily-MEAN heat index — doesn't heat illness come from the afternoon peak?",
     "The mean targets persistent day-and-night apparent heat (incl. warm nights). Trade-off: it under-weights the peak and is a DAILY PROXY, not hourly — so this is exposure CLASSIFICATION, not worker dose."),
    ("Heat index is defined on hourly data — how valid is a daily proxy?",
     "It's an approximate NWS proxy from GHCN-Daily temp + gridMET humidity. Regional gradient robust; single-county values & trend slopes are descriptive only."),
])
qa_slide("Likely questions — data quality & persistence", [
    ("With ~13% of county-days imputed and 22 counties fully imputed, can I trust the maps?",
     "Regionally yes, pixel-by-pixel no. Fully-imputed counties are interpolated, not observed — flag/exclude them for single-county claims. Statewide totals don't hinge on any one county."),
    ("Is a 2-day exceedance really a 'wave', and is a 48-day event credible?",
     "The >=2-day rule is a minimum-persistence filter; 44–53% of events are exactly 2 days, so 'elevated exposure' is fairer than 'wave' for much of it. Long tails are genuine persistent anomalies — lead with the median, not the max."),
    ("Some adjacent counties differ sharply — is that real climate?",
     "Largely no — it reflects the changing multi-station composite + IDW imputation, not micro-climate. That's why single-county rankings are presented cautiously."),
])
qa_slide("Likely questions — thresholds, windows & claims", [
    ("Why 85th and 95th specifically — aren't those arbitrary?",
     "Conventional relative heat thresholds (95th has published precedent), chosen to bracket a moderate and a severe cut. A fuller threshold sweep is a reasonable extension."),
    ("A percentile fixes the exceedance rate — isn't the day count preordained?",
     "The single-day rate is anchored (~15% / ~5%), but totals also depend on the >=2-day rule and the walk-forward baseline. The informative content is WHERE and WHEN heat clusters, not the grand total."),
    ("How sensitive are results to the 15-day window vs a calendar month?",
     "Not sensitive — the two windows agree at r≈0.99 per county. The consequential levers are the percentile and the floor/season choice."),
    ("What can these results actually support?",
     "County-level environmental apparent-heat EXPOSURE classification — not worker dose, causal injury, or official-advisory equivalence."),
])

# ---- Caveats ----
section("Caveats to keep on any results slide", RED)
bullets_slide("Caveats",
    [(0, "Daily heat-index PROXY (not hourly); temperature and humidity have different spatial support.", False),
     (0, "County temperature is a changing multi-station composite + IDW imputation -> single-county values noisy, regional gradient robust.", False),
     (0, "Year-round RELATIVE construct = persistent apparent-heat anomaly; ~63% of days are cool-season anomalies unless the >=80°F floor is applied.", False),
     (0, "NWS proxy is approximate and NOT an official advisory; trend slopes are descriptive only.", False),
     (0, "Exposure classification — NOT an injury or worker-heat-dose measure.", False)],
    note="Data: NOAA GHCN-Daily + gridMET + Census shapefile. Code (state-agnostic): github.com/eaadeniyi/gulf_heatwave")

import shutil
SCRATCH = r"C:\Users\eadeni1\AppData\Local\Temp\claude\C--Users-eadeni1-OneDrive---Louisiana-State-University-Documents-doc-heatWaveUS\2a61cee3-4efc-49b7-ace6-943ccd012cea\scratchpad\Heatwave_Results_Deck.pptx"
prs.save(SCRATCH)
print("[done] wrote", SCRATCH, "with", len(prs.slides._sldIdLst), "slides")
try:
    shutil.copyfile(SCRATCH, OUTPPTX)
    print("[copied] into repo:", OUTPPTX)
except PermissionError:
    print("[warn] repo copy locked; scratchpad copy is the current deck")
